"""Advanced file importers for PDF, DOCX, and other formats."""

import io
import re
from pathlib import Path
from typing import List, Optional


def import_pdf(content: bytes) -> List[dict]:
    """
    Import PDF file and extract text.

    Uses pypdf for basic extraction, marker-pdf for advanced layout analysis.
    """
    examples = []

    try:
        # Try marker-pdf first (better quality)
        from marker.convert import convert_single_pdf
        from marker.models import load_all_models

        models = load_all_models()
        full_text, images, out_meta = convert_single_pdf(
            io.BytesIO(content),
            models,
            max_pages=None,
            parallel_factor=1,
        )

        # Split by sections
        sections = _split_markdown_sections(full_text)
        for section in sections:
            if section.strip():
                examples.append({
                    "instruction": "Summarize or analyze the following text:",
                    "input": section.strip(),
                    "output": "",
                })

    except ImportError:
        # Fall back to pypdf
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(content))
            full_text = ""

            for page in reader.pages:
                text = page.extract_text()
                if text:
                    full_text += text + "\n\n"

            # Split into chunks
            chunks = _split_text_chunks(full_text, max_chars=4000)
            for chunk in chunks:
                if chunk.strip():
                    examples.append({
                        "instruction": "Summarize or analyze the following text:",
                        "input": chunk.strip(),
                        "output": "",
                    })

        except ImportError:
            raise ImportError("PDF support requires pypdf or marker-pdf. Install with: pip install pypdf")

    return examples


def import_docx(content: bytes) -> List[dict]:
    """
    Import DOCX file and extract structured content.
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("DOCX support requires python-docx. Install with: pip install python-docx")

    doc = Document(io.BytesIO(content))
    examples = []

    # Extract paragraphs
    current_section = []
    current_heading = ""

    for para in doc.paragraphs:
        text = para.text.strip()

        if not text:
            continue

        # Check if it's a heading
        if para.style.name.startswith("Heading"):
            # Save previous section
            if current_section:
                section_text = "\n\n".join(current_section)
                if section_text.strip():
                    examples.append({
                        "instruction": f"Explain the following section: {current_heading}" if current_heading else "Summarize the following text:",
                        "input": section_text.strip(),
                        "output": "",
                    })

            current_heading = text
            current_section = []
        else:
            current_section.append(text)

    # Don't forget the last section
    if current_section:
        section_text = "\n\n".join(current_section)
        if section_text.strip():
            examples.append({
                "instruction": f"Explain the following section: {current_heading}" if current_heading else "Summarize the following text:",
                "input": section_text.strip(),
                "output": "",
            })

    # Also extract tables
    for table in doc.tables:
        table_text = _extract_table_text(table)
        if table_text.strip():
            examples.append({
                "instruction": "Analyze the data in the following table:",
                "input": table_text,
                "output": "",
            })

    return examples


def import_pptx(content: bytes) -> List[dict]:
    """
    Import PowerPoint file and extract content from slides.

    Extracts:
    - Slide titles
    - Text from text boxes and shapes
    - Speaker notes
    - Tables
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError("PowerPoint support requires python-pptx. Install with: pip install python-pptx")

    prs = Presentation(io.BytesIO(content))
    examples = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        slide_title = ""

        # Extract shapes
        for shape in slide.shapes:
            # Get title
            if shape.has_text_frame:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title
                    slide_title = shape.text.strip()
                else:
                    text = shape.text.strip()
                    if text:
                        slide_content.append(text)

            # Extract tables
            if shape.has_table:
                table_text = _extract_pptx_table(shape.table)
                if table_text:
                    slide_content.append(f"Table:\n{table_text}")

        # Get speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()

        # Create example from slide content
        content_text = "\n\n".join(slide_content)

        if content_text.strip():
            instruction = f"Slide {slide_num}"
            if slide_title:
                instruction = f"Slide {slide_num}: {slide_title}"

            example = {
                "instruction": f"Explain the content from {instruction}:",
                "input": content_text,
                "output": "",
                "metadata": {
                    "source_type": "pptx",
                    "slide_number": slide_num,
                    "slide_title": slide_title,
                }
            }

            # Include speaker notes as additional context
            if notes_text:
                example["metadata"]["speaker_notes"] = notes_text

            examples.append(example)

    return examples


def _extract_pptx_table(table) -> str:
    """Extract text from a PPTX table."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(" | ".join(cells))
    return "\n".join(rows)


def import_xlsx(content: bytes) -> List[dict]:
    """
    Import Excel file and convert to dataset format.
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("Excel support requires openpyxl. Install with: pip install openpyxl")

    wb = load_workbook(io.BytesIO(content), data_only=True)
    examples = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # Get headers from first row
        headers = []
        for cell in sheet[1]:
            headers.append(str(cell.value) if cell.value else f"Column_{cell.column}")

        # Process rows
        for row_idx, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), start=2):
            row_data = {}
            for col_idx, value in enumerate(row):
                if col_idx < len(headers) and value is not None:
                    row_data[headers[col_idx]] = str(value)

            if row_data:
                # Convert to Q&A format if possible
                if _looks_like_qa(headers):
                    q_col = _find_column(headers, ["question", "query", "prompt", "input"])
                    a_col = _find_column(headers, ["answer", "response", "output", "reply"])

                    if q_col and a_col:
                        examples.append({
                            "instruction": row_data.get(q_col, ""),
                            "output": row_data.get(a_col, ""),
                        })
                        continue

                # Default: create structured data example
                data_str = "\n".join(f"{k}: {v}" for k, v in row_data.items())
                examples.append({
                    "instruction": f"Based on the following data from {sheet_name}:",
                    "input": data_str,
                    "output": "",
                })

    return examples


def import_html(content: bytes) -> List[dict]:
    """
    Import HTML file and extract article content.
    """
    try:
        from trafilatura import extract
    except ImportError:
        # Fall back to basic parsing
        return _basic_html_import(content)

    # Use trafilatura for high-quality extraction
    text = extract(content.decode("utf-8", errors="replace"))

    if not text:
        return _basic_html_import(content)

    # Split into chunks
    chunks = _split_text_chunks(text, max_chars=4000)
    examples = []

    for chunk in chunks:
        if chunk.strip():
            examples.append({
                "instruction": "Summarize or analyze the following content:",
                "input": chunk.strip(),
                "output": "",
            })

    return examples


def _basic_html_import(content: bytes) -> List[dict]:
    """Basic HTML import without trafilatura."""
    import html
    from html.parser import HTMLParser

    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
            self.skip_tags = {"script", "style", "noscript", "header", "footer", "nav"}
            self.current_skip = False

        def handle_starttag(self, tag, attrs):
            if tag in self.skip_tags:
                self.current_skip = True

        def handle_endtag(self, tag):
            if tag in self.skip_tags:
                self.current_skip = False

        def handle_data(self, data):
            if not self.current_skip:
                text = data.strip()
                if text:
                    self.text.append(text)

    parser = TextExtractor()
    parser.feed(content.decode("utf-8", errors="replace"))

    full_text = " ".join(parser.text)
    chunks = _split_text_chunks(full_text, max_chars=4000)

    return [
        {"instruction": "Summarize the following:", "input": chunk, "output": ""}
        for chunk in chunks
        if chunk.strip()
    ]


def import_url(url: str) -> List[dict]:
    """
    Import content from a URL.
    """
    import httpx

    response = httpx.get(url, follow_redirects=True, timeout=30)
    response.raise_for_status()

    content_type = response.headers.get("content-type", "")

    if "application/json" in content_type:
        # JSON content
        from dataforge_core.importers import _import_json
        return _import_json(response.content)

    elif "text/html" in content_type:
        return import_html(response.content)

    elif "text/plain" in content_type:
        from dataforge_core.importers import _import_txt
        return _import_txt(response.content)

    else:
        # Try as text
        return [{"text": response.text}]


def import_huggingface_dataset(dataset_path: str, split: str = "train", max_examples: int = 10000) -> List[dict]:
    """
    Import dataset from HuggingFace Hub.

    Args:
        dataset_path: HuggingFace dataset identifier (e.g., "tatsu-lab/alpaca")
        split: Dataset split to load
        max_examples: Maximum number of examples to import
    """
    try:
        from datasets import load_dataset
    except ImportError:
        raise ImportError("HuggingFace support requires datasets library. Install with: pip install datasets")

    dataset = load_dataset(dataset_path, split=split, streaming=True)

    examples = []
    for i, item in enumerate(dataset):
        if i >= max_examples:
            break
        examples.append(dict(item))

    return examples


# =============================================================================
# Helper Functions
# =============================================================================

def _split_markdown_sections(text: str) -> List[str]:
    """Split markdown text by headers."""
    # Split by ## or # headers
    sections = re.split(r'\n#{1,2}\s+', text)
    return [s.strip() for s in sections if s.strip()]


def _split_text_chunks(text: str, max_chars: int = 4000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks."""
    if len(text) <= max_chars:
        return [text]

    chunks = []
    start = 0

    while start < len(text):
        end = start + max_chars

        # Try to break at a paragraph or sentence boundary
        if end < len(text):
            # Look for paragraph break
            para_break = text.rfind("\n\n", start, end)
            if para_break > start + max_chars // 2:
                end = para_break

            # Or sentence break
            else:
                sentence_break = text.rfind(". ", start, end)
                if sentence_break > start + max_chars // 2:
                    end = sentence_break + 1

        chunks.append(text[start:end].strip())
        start = end - overlap

    return chunks


def _extract_table_text(table) -> str:
    """Extract text from a DOCX table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))
    return "\n".join(rows)


def _looks_like_qa(headers: List[str]) -> bool:
    """Check if headers suggest a Q&A format."""
    headers_lower = [h.lower() for h in headers]
    qa_patterns = [
        ("question", "answer"),
        ("query", "response"),
        ("prompt", "completion"),
        ("input", "output"),
    ]

    for q, a in qa_patterns:
        if any(q in h for h in headers_lower) and any(a in h for h in headers_lower):
            return True

    return False


def _find_column(headers: List[str], patterns: List[str]) -> Optional[str]:
    """Find a column header matching any of the patterns."""
    headers_lower = {h.lower(): h for h in headers}

    for pattern in patterns:
        for header_lower, header in headers_lower.items():
            if pattern in header_lower:
                return header

    return None
